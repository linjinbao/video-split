import subprocess
import os
import tempfile
from pptx import Presentation
from PIL import Image
import shutil

# 配置参数
VIDEO_OUTPUT_DIR = "./video_output_1734595582547-8f606097a94a35b19e6590cc6e07de437d6cb8b9-video_cb208bb0"  # 分割视频的输出目录
PPT_FILE = "ppt.pptx"  # 输入的 PPT 文件路径
FINAL_VIDEO = "final_output.mp4"  # 最终合成的视频文件路径

# 1. 将 PPT 文件转换为图片
def ppt_to_images(ppt_file, output_dir):
    os.makedirs(output_dir, exist_ok=True)

    # 打开 PPT 文件
    prs = Presentation(ppt_file)

    # 遍历 PPT 每一页
    for i, slide in enumerate(prs.slides):
        # 创建图片文件路径
        image_path = os.path.join(output_dir, f"slide_{i + 1}.png")
        
        # 暂时假设这里用一些方法将 PPT 页转换为图片
        # 使用 python-pptx 转图片，这里需要外部库，如 unoconv 或 libreoffice
        slide.shapes._spTree
        # 需要你实现 PPT 转图片的代码
        slide.shapes[0].element

        # 示例转换图片代码，通常需要转化为图片再保存
        img = Image.new("RGB", (1920, 1080))  # 假设背景是 1920x1080
        img.save(image_path)
        print(f"已将 PPT 页 {i + 1} 转换为图片 {image_path}")
    
    return len(prs.slides)  # 返回 PPT 页数

# 2. 将视频片段和 PPT 图片合成成视频
def merge_video_with_image(video_dir, image_dir, final_video_path):
    video_files = sorted([f for f in os.listdir(video_dir) if f.endswith(".mp4")])
    image_files = sorted([f for f in os.listdir(image_dir) if f.endswith(".png")])

    # 确保视频和图片数量一致
    if len(video_files) != len(image_files):
        print("视频和图片数量不匹配！")
        return

    # 临时合成的视频片段列表
    temp_files = []

    for i in range(len(video_files)):
        video_file = os.path.join(video_dir, video_files[i])
        image_file = os.path.join(image_dir, image_files[i])

        # 合成命令：将视频缩小到 240x135 并放到右下角
        temp_output = os.path.join(temp_dir, f"temp_output_{i}.mp4")
        temp_files.append(temp_output)

        # 使用 ffmpeg 合成视频和图片
        command = [
            "ffmpeg", "-i", video_file, "-i", image_file, "-filter_complex", 
            "[0:v]scale=240:135[small];[1:v][small]overlay=W-w-10:H-h-10", 
            "-c:a", "aac", "-strict", "experimental", "-y", temp_output
        ]
        subprocess.run(command)
        print(f"已合成视频片段 {i + 1}：{temp_output}")

    # 合并所有片段
    with open(os.path.join(temp_dir, "temp_list.txt"), "w") as f:
        for temp_file in temp_files:
            f.write(f"file '{temp_file}'\n")

    # 使用 ffmpeg 合并视频片段
    subprocess.run(["ffmpeg", "-f", "concat", "-safe", "0", "-i", os.path.join(temp_dir, "temp_list.txt"), "-c", "copy", final_video_path])
    print(f"视频合成完成，最终视频：{final_video_path}")

    # 清理临时文件
    for temp_file in temp_files:
        os.remove(temp_file)
    os.remove(os.path.join(temp_dir, "temp_list.txt"))

# 主程序
if __name__ == "__main__":
    # 创建临时目录
    with tempfile.TemporaryDirectory() as temp_dir:
        print(f"临时目录创建在：{temp_dir}")

        # 1. 将 PPT 转换为图片
        images_dir = os.path.join(temp_dir, "ppt_images")
        num_pages = ppt_to_images(PPT_FILE, images_dir)

        # 2. 合成视频
        final_video_path = os.path.join(temp_dir, FINAL_VIDEO)
        merge_video_with_image(VIDEO_OUTPUT_DIR, images_dir, final_video_path)

        # 移动最终视频到当前工作目录
        shutil.move(final_video_path, FINAL_VIDEO)
        print(f"最终视频已保存为：{FINAL_VIDEO}")

